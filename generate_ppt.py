from pptx import Presentation
from pptx.util import Inches, Pt

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]

    title_placeholder.text = title

    tf = body_placeholder.text_frame
    for point in content:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)

def main():
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Rigup Төсөл"
    subtitle.text = "React frontend + SOA backend (AI агент) + Docker эхлүүлэгч\nТөслийн танилцуулга"

    # Slide 2
    add_slide(
        prs,
        "Төслийн Архитектур",
        [
            "• Frontend (React/CRA): Хэрэглэгчийн интерфэйс.",
            "• Core-API (Node.js/Express): Үндсэн HTTP API систем.",
            "• Agent-Worker (Node.js/OpenAI SDK): Цаана ажиллах хиймэл оюун ухааны даалгавар боловсруулагч.",
            "• MySQL (MySQL 8): Мэдээллийн сан."
        ]
    )

    # Slide 3
    add_slide(
        prs,
        "Бүрэлдэхүүн хэсгүүдийн үүрэг",
        [
            "• Frontend: Газрын зураг, кафе, PC сүлжээ, захиалга болон чат зэргийг харуулна.",
            "• Core-API: Хэрэглэгч нэвтрэх, кафе/PC мэдээлэл, захиалга, төлбөр болон AI даалгаврыг бүртгэх үүрэгтэй.",
            "• Agent-Worker: 'agent_tasks' хүснэгтээс даалгавар уншиж, LLM ашиглан хариуг боловсруулж буцааж бичнэ.",
            "• MySQL: Бүх мэдээллийг (кафе, pc, захиалга, агентийн дараалал) нэг дор хадгална."
        ]
    )

    # Slide 4
    add_slide(
        prs,
        "AI Агентийн ажиллах урсгал",
        [
            "1. Хөтчөөс 'POST /api/agent/chat' руу хэрэглэгчийн зурвас илгээгдэнэ.",
            "2. 'core-api' нь баталгаажуулж, даалгаврыг 'agent_tasks' руу бүртгэж хариуг хүлээдэг.",
            "3. 'agent-worker' нь даалгаврыг хүлээн авч OpenAI/Groq ашиглан боловсруулаад хариуг буцаадаг.",
            "4. 'core-api' хариуг хүлээн авч хэрэглэгч рүү илгээнэ."
        ]
    )

    # Slide 5
    add_slide(
        prs,
        "Хурдан эхлүүлэх заавар (Quick Start)",
        [
            "• Кодыг татах: git clone <repo> && cd Rigup",
            "• Тохиргоог хуулах: cp .env.example .env (JWT_SECRET, MYSQL_PASSWORD, AI_API_KEY өөрчлөх)",
            "• Төслийг ажиллуулах: npm start (эсвэл ./start.sh, make start)",
            "• Frontend хаяг: http://localhost:5173",
            "• API хаяг: http://localhost:5500/api",
            "• MySQL: localhost:33306"
        ]
    )

    # Slide 6
    add_slide(
        prs,
        "Үндсэн командууд болон тохиргоо",
        [
            "• npm stop: Төслийг зогсоох",
            "• npm run logs: Бүх сервисийн логийг харах",
            "• npm run rebuild: Дахин компайл хийж ажиллуулах",
            "• npm run clean: Бүх өгөгдлийг устгаж цэвэрлэх",
            "• AI-г ажиллуулахын тулд .env файлд 'AI_API_KEY'-г заавал нэмэх шаардлагатай."
        ]
    )

    prs.save('Rigup_Presentation_MN.pptx')
    print("PowerPoint presentation generated successfully!")

if __name__ == '__main__':
    main()
